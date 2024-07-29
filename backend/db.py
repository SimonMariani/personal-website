from langchain.text_splitter import RecursiveCharacterTextSplitter
from pymilvus import MilvusClient, DataType
from openai import OpenAI
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import os

try:
    from dotenv import load_dotenv
    load_dotenv(override=True)
except:
    pass

# Load the environment variables
db_url = os.environ['DB_URL']
db_token = os.environ['DB_TOKEN']

# Connect to clients and instatiate the splitter and embedding function
milvus_client  = MilvusClient(uri=db_url, token=db_token)
openai_client = OpenAI(api_key=os.environ['OPENAIKEY'])
splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)

def emb_text(text):
    """
    Embedding function that returns an embedding for a given text.
    """

    return (
        openai_client.embeddings.create(input=text, model="text-embedding-3-small")
        .data[0]
        .embedding
    )

# Create the schema given the embedding dimension
embedding_dim = len(emb_text("text to get the embedding dimension"))

schema = milvus_client.create_schema(
    auto_id=True,
    enable_dynamic_field=True,
)

schema.add_field(field_name="id", datatype=DataType.INT64, is_primary=True)
schema.add_field(field_name="filename", datatype=DataType.VARCHAR, max_length=256)
schema.add_field(field_name="vector", datatype=DataType.FLOAT_VECTOR, dim=embedding_dim)
schema.add_field(field_name="text", datatype=DataType.VARCHAR, max_length=65535)

# Create the index params
index_params = milvus_client.prepare_index_params()

index_params.add_index(
    field_name="id",
    index_type="STL_SORT"
)

index_params.add_index(
    field_name="vector", 
    index_type="IVF_FLAT",
    metric_type="IP",
    params={ "nlist": 128 }
)



def get_file_contens(file_path):
    """
    Read the content of the file and return the content as a string the content splitted in chunks.
    """

    file_extension = os.path.splitext(file_path)[1]

    if file_extension == '.pdf':
        pdf_reader = PdfReader(file_path)
        file_content = '\n'.join([page.extract_text() for page in pdf_reader.pages])

    elif file_extension == '.docx':
        doc = Document(file_path)
        file_content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])

    elif file_extension == '.pptx':
        prs = Presentation(file_path)
        file_content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame])

    elif file_extension == '.csv':
        df = pd.read_csv(file_path, low_memory=False)
        file_content = df.to_string(index=False)

    elif file_extension in ['.xlsx', '.xls']:
        df = pd.read_excel(file_path)
        file_content = df.to_string(index=False)

    else:
        f = open(file_path, "r")
        file_content = f.read()

    file_content_splitted = splitter.split_text(file_content)

    return file_content, file_content_splitted


def add_text_array_to_vector_db(texts, file_name, collection_name):
    """
    Add an array of strings to a given collection in the vector database.
    """

    # Since we do not have a system of adding new files or removing old files we simply just start over every time.
    if milvus_client.has_collection(collection_name):
        milvus_client.drop_collection(collection_name)

    milvus_client.create_collection(
        collection_name=collection_name,
        dimension=embedding_dim,
        consistency_level="Strong",  # Strong consistency level
        schema=schema,
        index_params=index_params
    )

    # Create all the new entries
    embeddings = [emb_text(text) for text in texts]

    data = [{
        "filename": file_name,
        "vector": embedding,
        "text": text,
    } for embedding, text in zip(embeddings, texts)]

    milvus_client.insert(collection_name=collection_name, data=data)


def remove_from_vector_db(self):
    """
    Remove all the entries from the vector database that are related to this document.
    """
    
    res = milvus_client.delete(
        collection_name=self.collection_name,
        filter=f"id_document == '{str(self.id)}'"
    )

    res_splitted = milvus_client.delete(
        collection_name=self.collection_name_splitted,
        filter=f"id_document == '{str(self.id)}'"
    )
    # TODO handle errors


def retrieve_relevant_content(query, use_splitted=True, limit=100, truncate_text=None):
    """
    Retrieve the relevant content from the Milvus database based on the query.
    """

    # Either retrieve from the splitted collection or from the collection with the full documents
    collection = "documents_splitted" if use_splitted else "documents" 

    texts = []
    if milvus_client.has_collection(collection) and query is not None:

        search_res = milvus_client.search(
            collection_name=collection,
            data=[
                emb_text(query)
            ],
            limit=limit,
            search_params={"metric_type": "IP", "params": {}},  # Inner product distance
            output_fields=["text", "id", "filename"],  # Return the text field
        )

        texts = [res["entity"]["text"][:truncate_text] for res in search_res[0]]

    return texts


if __name__ == "__main__":
    # When ran as a script we will process all the files in the files directory.
    print("running db.py as script")

    directory = 'files'
    for file in os.listdir(directory):
        print("Processing file: ", file)

        file_path = os.path.join(directory, file)
        file_name = os.path.basename(file_path)

        file_content, file_content_splitted = get_file_contens(file_path)
        add_text_array_to_vector_db([file_content], file_name, 'documents')
        add_text_array_to_vector_db(file_content_splitted, file_name, 'documents_splitted')





 # Create the collection if it does not exist
    # if not milvus_client.has_collection(collection_name):
    #     milvus_client.create_collection(
    #         collection_name=collection_name,
    #         dimension=embedding_dim,
    #         consistency_level="Strong",  # Strong consistency level
    #         schema=schema,
    #         index_params=index_params
    #     )                
    
    # Delete all the old entries from this document
    # milvus_client.delete(
    #     collection_name=collection_name,
    #     filter=f"id_document == '{str(self.id)}'"
    # )