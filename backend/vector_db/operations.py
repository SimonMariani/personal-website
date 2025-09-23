from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import os
from .connections import milvus_client, openai_client, embedding_model
from .utility import emb_text
from .setup import embedding_dim, schema, index_params


##############
### DB ADD ###
##############


def add_to_vector_db(collection_name, filename, file_path):
    """
    Add all the relevant information for this document to the vector database.
    """

    file_content = get_file_content(file_path)
    file_content_chunks = split_text(file_content)
    add_chunks(file_content_chunks, collection_name, filename, file_path)


def get_file_content(file_path):
    """
    Read the content of the file and return the content as a string the content splitted in chunks.
    """

    # Obtain the file extension
    file_extension = os.path.splitext(file_path)[1]

    if file_extension == ".pdf":
        pdf_reader = PdfReader(file_path)
        file_content = [(page.extract_text(), {"page": i}) for i, page in enumerate(pdf_reader.pages)]

    elif file_extension == ".docx":
        doc = Document(file_path)
        file_content = [(para.text, {"page": i}) for i, para in enumerate(doc.paragraphs) if para.text.strip()]

    elif file_extension == ".pptx":
        prs = Presentation(file_path)
        file_content = []
        for i, slide in enumerate(prs.slides):
            slide_text = "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip())
            if slide_text:
                file_content.append((slide_text, {"page": i}))

    elif file_extension == ".csv":
        df = pd.read_csv(file_path, low_memory=False)
        file_content = [(row, {"page": i}) for i, row in enumerate(df.astype(str).apply(lambda row: ", ".join(row), axis=1))]

    elif file_extension in [".xlsx", ".xls"]:
        df = pd.read_excel(file_path)
        file_content = [(row, {"page": i}) for i, row in enumerate(df.astype(str).apply(lambda row: ", ".join(row), axis=1))]

    else:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()
            file_content = [(line.strip(), {"page": i}) for i, line in enumerate(lines) if line.strip()]

    return file_content


def split_text(file_content, chunk_size=1000, chunk_overlap=50):
    """
    Split text into chunks of a given size with a given overlap.
    """

    all_splits = []
    for text, meta in file_content:
        # Initialize variables
        start = 0
        splits = []

        # Loop through the content to create chunks
        while start < len(text):
            end = min(start + chunk_size, len(text))
            splits.append((text[start:end], meta))

            # Move the starting point forward by chunk_size - chunk_overlap
            start += chunk_size - chunk_overlap

        # Add the splits to the all_splits list
        all_splits.extend(splits)

    return all_splits


def add_chunks(file_content_chunks, collection_name, filename, file_path):
    """
    Add an array of strings to a given collection in the vector database.
    """

    # The name of the collection is the company id
    if not milvus_client.has_collection(collection_name):
        milvus_client.create_collection(
            collection_name=collection_name,
            dimension=embedding_dim,
            consistency_level="Strong",
            schema=schema,
            index_params=index_params,
        )

    # Delete all the old entries from this document
    milvus_client.delete(collection_name=collection_name, filter=f"filename == '{filename}'")

    # Create embeddings for all the chunks
    embeddings = [(emb_text(text, openai_client, embedding_model), text, meta) for text, meta in file_content_chunks]

    # Add the embeddings to the database
    data = [
        {
            "filename": filename,
            "file_path": file_path,
            "vector": embedding,
            "text": text,
            "page": meta["page"],
        }
        for embedding, text, meta in embeddings
    ]
    milvus_client.insert(collection_name=collection_name, data=data)


####################
### BULK UPDATE ###
####################


def update_vector_db(collection_name, directory, overwrite=False):
    """
    Update the vector database from the files in the given directory, this consists of adding the files and removing the files that no longer exist.
    """

    print("Starting vector database update...")

    add_files_to_vector_db(collection_name, directory=directory, overwrite=overwrite)
    remove_deleted_files_from_vector_db(collection_name, directory=directory)

    print("Vector database update complete.")


def add_files_to_vector_db(collection_name, directory, overwrite=False):
    """
    Update the vector database from the files in the files directory.
    """

    print("Adding files to vector database...")

    for file in os.listdir(directory):

        # If the file already exists in the vector database and overwrite is False, we skip it
        if filename_exists_in_vector_db(collection_name, file) and not overwrite:
            print(f"File {file} already exists in the vector database. Skipping.")
            continue

        # Else we process it
        print("Processing file: ", file)
        add_to_vector_db(collection_name, file, os.path.join(directory, file))

    print("Finished adding files to vector database.")


def remove_deleted_files_from_vector_db(collection_name, directory="documents"):
    """
    Ensure the Milvus collection only contains entries for files that still exist in the given directory. Deletes entries for missing files.
    """

    print("Removing deleted files from vector database...")

    existing_files = set(os.listdir(directory))

    # Ensure the collection exists
    if not milvus_client.has_collection(collection_name):
        print(f"Collection '{collection_name}' does not exist.")
        return

    # Get all filenames stored in Milvus
    results = milvus_client.query(
        collection_name=collection_name, filter="", output_fields=["filename"], limit=10000  # pick a safe upper bound for your dataset
    )
    db_files = {res["filename"] for res in results}

    # Find DB entries that no longer have a file on disk
    deleted_files = db_files - existing_files

    if not deleted_files:
        print("No files to remove. Milvus is in sync with the directory.")
        return

    # Remove each missing file from Milvus
    for filename in deleted_files:
        print(f"Removing {filename} from vector database (file no longer exists).")
        milvus_client.delete(collection_name=collection_name, filter=f'filename == "{filename}"')


###################
### DB RETRIEVE ###
###################


def filename_exists_in_vector_db(collection_name, filename):
    """
    Check if a filename exists in the vector database collection. Returns True if exists, False otherwise.
    """

    # If the collection does not exist, return False
    if not milvus_client.has_collection(collection_name):
        return False

    # Query the collection for the filename
    results = milvus_client.query(collection_name=collection_name, filter=f"filename == '{filename}'", output_fields=["filename"], limit=1)

    # Return True if results are found, False otherwise
    return len(results) > 0


def retrieve_from_query(query, collection_name, limit=100):

    # Validation checks
    if query is None:
        raise ValueError("No query provided.")

    if not milvus_client.has_collection(collection_name):
        raise ValueError("The collection does not exist.")

    # Search the database
    search_res = milvus_client.search(
        collection_name=collection_name,
        data=[emb_text(query, openai_client, embedding_model)],
        limit=limit,
        search_params={"metric_type": "IP", "params": {}},  # Inner product distance
        output_fields=["id", "filename", "file_path", "text", "page"],  # Return the text field
    )[0]

    # Because Milvus introduces some search response object we need to convert it to a python list with dicts
    transformed_search_res = []
    for res in search_res:
        transformed_response = {}
        for key in res:
            transformed_response[key] = res[key]
        transformed_search_res.append(transformed_response)

    # Finally we return the transforrmed search results
    return transformed_search_res


#################
### DB REMOVE ###
#################


def remove_collection_from_vector_db(collection_name):
    """
    Remove an entire collection from the vector database.
    """
    if milvus_client.has_collection(collection_name):
        milvus_client.drop_collection(collection_name)
        print(f"Collection '{collection_name}' has been removed from the vector database.")
    else:
        print(f"Collection '{collection_name}' does not exist.")
