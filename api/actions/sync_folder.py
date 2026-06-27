import os
from psycopg import sql
from pgvector.psycopg import Vector
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
from connections.db_connection import pool
from connections.model_connection import client
from config import EMBEDDING_MODEL, EMBEDDING_DIM, CHUNK_SIZE, CHUNK_OVERLAP, SEARCH_COLLECTION

# Max inputs per OpenAI embeddings request; the API rejects larger batches.
EMBED_BATCH_SIZE = 1000

##################
### OPERATIONS ###
##################


def sync_directory(directory, collection=SEARCH_COLLECTION, overwrite=False, only_file=None):
    """
    Match a collection to a directory: upload new/updated files and drop removed ones.
    """
    # Either sync a single requested file or every file in the directory.
    files = [only_file] if only_file else os.listdir(directory)
    stored = _stored_filenames(collection)

    # Upload each file, skipping ones already stored unless overwrite is set.
    for filename in files:
        if filename in stored and not overwrite:
            print(f"Skipping '{filename}' (already stored).")
            continue
        print(f"Uploading '{filename}'...")
        upload_file(filename, os.path.join(directory, filename), collection)

    # Drop rows for files that no longer exist on disk.
    for filename in stored - set(os.listdir(directory)):
        print(f"Removing '{filename}' (file no longer exists).")
        with pool.connection() as conn:
            conn.execute(sql.SQL("DELETE FROM {} WHERE filename = %s").format(sql.Identifier(collection)), (filename,))


def upload_file(filename, file_path, collection=SEARCH_COLLECTION):
    """
    Read, chunk, embed and store a file, replacing any previous version of it.
    """
    # Split the file into text chunks; skip files that produced no text.
    chunks = _load_chunks(file_path)
    if not chunks:
        return

    # Embed every chunk and pair the vectors back up with their text/page.
    vectors = _embed([text for text, _ in chunks])
    rows = [(filename, file_path, page, text, Vector(vector)) for (text, page), vector in zip(chunks, vectors)]

    with pool.connection() as conn, conn.cursor() as cur:
        # Make sure the collection exists, then replace this file's rows.
        _create_table(conn, collection)
        cur.execute(sql.SQL("DELETE FROM {} WHERE filename = %s").format(sql.Identifier(collection)), (filename,))
        cur.executemany(
            sql.SQL("INSERT INTO {} (filename, file_path, page, text, vector) VALUES (%s, %s, %s, %s, %s)").format(sql.Identifier(collection)),
            rows,
        )


#################
### EMBEDDING ###
#################


def _embed(texts):
    """
    Embed a list of texts, batching requests to stay within the API's input limit.
    """
    vectors = []
    for start in range(0, len(texts), EMBED_BATCH_SIZE):
        batch = texts[start : start + EMBED_BATCH_SIZE]
        # Sort each response back into input order before appending.
        response = client.embeddings.create(input=batch, model=EMBEDDING_MODEL)
        vectors.extend(item.embedding for item in sorted(response.data, key=lambda item: item.index))
    return vectors


########################
### DOCUMENT LOADING ###
########################


def _load_chunks(file_path):
    """
    Read a file and return a list of (text, page) chunks.
    """
    step = CHUNK_SIZE - CHUNK_OVERLAP
    chunks = []

    # Walk each page/section of the file and slice it into overlapping chunks.
    for text, page in _read_pages(file_path):
        # Coalesce missing text and strip NUL bytes, which Postgres text columns reject.
        text = (text or "").replace("\x00", "")
        for start in range(0, len(text), step):
            chunks.append((text[start : start + CHUNK_SIZE], page))

    return chunks


def _read_pages(file_path):
    """
    Read a file into a list of (text, page) tuples. "Page" is the page number for
    PDFs/slides and the paragraph/row index for the other formats.
    """
    extension = os.path.splitext(file_path)[1].lower()

    # PDF: one entry per page.
    if extension == ".pdf":
        return [(page.extract_text(), i) for i, page in enumerate(PdfReader(file_path).pages)]

    # Word: one entry per non-empty paragraph.
    if extension == ".docx":
        return [(p.text, i) for i, p in enumerate(Document(file_path).paragraphs) if p.text.strip()]

    # PowerPoint: one entry per slide, joining all of its text frames.
    if extension == ".pptx":
        pages = []
        for i, slide in enumerate(Presentation(file_path).slides):
            text = "\n".join(s.text for s in slide.shapes if s.has_text_frame and s.text.strip())
            if text:
                pages.append((text, i))
        return pages

    # Spreadsheets: one entry per row.
    if extension == ".csv":
        return _rows(pd.read_csv(file_path, low_memory=False))
    if extension in (".xlsx", ".xls"):
        return _rows(pd.read_excel(file_path))

    # Anything else: treat it as plain text, one entry per non-empty line.
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return [(line.strip(), i) for i, line in enumerate(f) if line.strip()]


def _rows(df):
    """
    Render each dataframe row as a comma-joined string paired with its row index.
    """
    rows = df.astype(str).apply(lambda row: ", ".join(row), axis=1)
    return [(row, i) for i, row in enumerate(rows)]


########################
### DATABASE HELPERS ###
########################


def _create_table(conn, collection):
    """
    Create the collection's table and its HNSW (cosine) index if they don't exist.
    """
    # One table per collection, with a fixed-size vector column.
    conn.execute(
        sql.SQL(
            "CREATE TABLE IF NOT EXISTS {table} (id BIGSERIAL PRIMARY KEY, "
            "filename VARCHAR(256), file_path VARCHAR(256), page BIGINT, text TEXT, vector vector({dim}))"
        ).format(table=sql.Identifier(collection), dim=sql.Literal(EMBEDDING_DIM))
    )
    # An HNSW index on the vector column makes cosine-distance search fast.
    conn.execute(
        sql.SQL("CREATE INDEX IF NOT EXISTS {index} ON {table} USING hnsw (vector vector_cosine_ops)").format(
            index=sql.Identifier(f"{collection}_vector_idx"), table=sql.Identifier(collection)
        )
    )


def _stored_filenames(collection):
    """
    Return the set of filenames already stored in a collection (empty if it has none).
    """
    with pool.connection() as conn:
        # A missing collection simply has no stored filenames.
        if conn.execute("SELECT to_regclass(%s)", (collection,)).fetchone()[0] is None:
            return set()
        rows = conn.execute(sql.SQL("SELECT DISTINCT filename FROM {}").format(sql.Identifier(collection))).fetchall()
    return {row[0] for row in rows}
