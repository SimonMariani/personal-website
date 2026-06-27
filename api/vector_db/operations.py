from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import os
from psycopg import sql
from pgvector.psycopg import Vector
from .connections import pg_pool, openai_client, embedding_model
from .utility import emb_text
from .setup import embedding_dim

#######################
### COLLECTION UTIL ###
#######################


def _collection_exists(conn, collection_name):
    """
    Check whether the table backing a collection exists.
    """
    row = conn.execute("SELECT to_regclass(%s)", (collection_name,)).fetchone()
    return row[0] is not None


def _ensure_collection(conn, collection_name):
    """
    Create the table and HNSW (cosine) index for a collection if they don't exist yet.
    """
    if _collection_exists(conn, collection_name):
        return

    conn.execute(
        sql.SQL(
            "CREATE TABLE IF NOT EXISTS {tbl} ("
            "id BIGSERIAL PRIMARY KEY, "
            "filename VARCHAR(256), "
            "file_path VARCHAR(256), "
            "page BIGINT, "
            "text TEXT, "
            "vector vector({dim}))"
        ).format(tbl=sql.Identifier(collection_name), dim=sql.Literal(embedding_dim))
    )
    conn.execute(
        sql.SQL("CREATE INDEX IF NOT EXISTS {idx} ON {tbl} USING hnsw (vector vector_cosine_ops)").format(
            idx=sql.Identifier(f"{collection_name}_vector_hnsw_idx"),
            tbl=sql.Identifier(collection_name),
        )
    )


##############
### DB ADD ###
##############


def add_to_vector_db(collection_name, filename, file_path):
    """
    Add all the relevant information for this document to the vector database.
    """

    file_content = get_file_content(file_path)
    file_content_chunks = split_text(file_content)
    add_chunks(file_content_chunks, collection_name, filename, file_path)


def get_file_content(file_path):
    """
    Read the content of the file and return the content as a string the content splitted in chunks.
    """

    # Obtain the file extension
    file_extension = os.path.splitext(file_path)[1]

    if file_extension == ".pdf":
        pdf_reader = PdfReader(file_path)
        file_content = [(page.extract_text(), {"page": i}) for i, page in enumerate(pdf_reader.pages)]

    elif file_extension == ".docx":
        doc = Document(file_path)
        file_content = [(para.text, {"page": i}) for i, para in enumerate(doc.paragraphs) if para.text.strip()]

    elif file_extension == ".pptx":
        prs = Presentation(file_path)
        file_content = []
        for i, slide in enumerate(prs.slides):
            slide_text = "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip())
            if slide_text:
                file_content.append((slide_text, {"page": i}))

    elif file_extension == ".csv":
        df = pd.read_csv(file_path, low_memory=False)
        file_content = [(row, {"page": i}) for i, row in enumerate(df.astype(str).apply(lambda row: ", ".join(row), axis=1))]

    elif file_extension in [".xlsx", ".xls"]:
        df = pd.read_excel(file_path)
        file_content = [(row, {"page": i}) for i, row in enumerate(df.astype(str).apply(lambda row: ", ".join(row), axis=1))]

    else:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()
            file_content = [(line.strip(), {"page": i}) for i, line in enumerate(lines) if line.strip()]

    return file_content


def split_text(file_content, chunk_size=1000, chunk_overlap=50):
    """
    Split text into chunks of a given size with a given overlap.
    """

    all_splits = []
    for text, meta in file_content:
        # Initialize variables
        start = 0
        splits = []

        # Loop through the content to create chunks
        while start < len(text):
            end = min(start + chunk_size, len(text))
            splits.append((text[start:end], meta))

            # Move the starting point forward by chunk_size - chunk_overlap
            start += chunk_size - chunk_overlap

        # Add the splits to the all_splits list
        all_splits.extend(splits)

    return all_splits


def add_chunks(file_content_chunks, collection_name, filename, file_path):
    """
    Add an array of strings to a given collection in the vector database.
    """

    # Create embeddings for all the chunks
    embeddings = [(emb_text(text, openai_client, embedding_model), text, meta) for text, meta in file_content_chunks]

    with pg_pool.connection() as conn:
        # The name of the collection is the company id
        _ensure_collection(conn, collection_name)

        # Delete all the old entries from this document
        conn.execute(
            sql.SQL("DELETE FROM {} WHERE filename = %s").format(sql.Identifier(collection_name)),
            (filename,),
        )

        # Add the embeddings to the database
        with conn.cursor() as cur:
            cur.executemany(
                sql.SQL("INSERT INTO {} (filename, file_path, page, text, vector) VALUES (%s, %s, %s, %s, %s)").format(sql.Identifier(collection_name)),
                [(filename, file_path, meta["page"], text, Vector(embedding)) for embedding, text, meta in embeddings],
            )


####################
### BULK UPDATE ###
####################


def update_vector_db(collection_name, directory, overwrite=False, file=None):
    """
    Update the vector database from the files in the given directory, this consists of adding the files and removing the files that no longer exist.
    """

    print(">>>Starting vector database update...")

    add_files_to_vector_db(collection_name, directory=directory, overwrite=overwrite, file=file)
    remove_deleted_files_from_vector_db(collection_name, directory=directory)

    print(">>>Vector database update complete.")


def add_files_to_vector_db(collection_name, directory, overwrite=False, file=None):
    """
    Update the vector database from the files in the files directory.
    """

    print(f'>>>Adding files from directory "{directory}" to vector database...')

    # Get the list of files in the directory or the single file if specified
    files = [file] if file else os.listdir(directory)
    for file in files:

        print(f"Checking file: {file}")

        # If the file already exists in the vector database and overwrite is False, we skip it
        if filename_exists_in_vector_db(collection_name, file) and not overwrite:
            print(f"File {file} already exists in the vector database. Skipping.")
            continue

        # Else we process it
        print("Processing file: ", file)
        add_to_vector_db(collection_name, file, os.path.join(directory, file))

    print(">>>Finished adding files to vector database.")


def remove_deleted_files_from_vector_db(collection_name, directory="documents"):
    """
    Ensure the collection only contains entries for files that still exist in the given directory. Deletes entries for missing files.
    """

    print(">>>Removing deleted files from vector database...")

    existing_files = set(os.listdir(directory))

    with pg_pool.connection() as conn:
        # Ensure the collection exists
        if not _collection_exists(conn, collection_name):
            print(f"Collection '{collection_name}' does not exist.")
            return

        # Get all filenames stored in the collection
        rows = conn.execute(sql.SQL("SELECT DISTINCT filename FROM {}").format(sql.Identifier(collection_name))).fetchall()
        db_files = {row[0] for row in rows}

        # Find DB entries that no longer have a file on disk
        deleted_files = db_files - existing_files

        if not deleted_files:
            print("No files to remove. The vector database is in sync with the directory.")
            return

        # Remove each missing file from the collection
        for filename in deleted_files:
            print(f"Removing {filename} from vector database (file no longer exists).")
            conn.execute(
                sql.SQL("DELETE FROM {} WHERE filename = %s").format(sql.Identifier(collection_name)),
                (filename,),
            )


###################
### DB RETRIEVE ###
###################


def filename_exists_in_vector_db(collection_name, filename):
    """
    Check if a filename exists in the vector database collection. Returns True if exists, False otherwise.
    """

    with pg_pool.connection() as conn:
        # If the collection does not exist, return False
        if not _collection_exists(conn, collection_name):
            return False

        # Query the collection for the filename
        row = conn.execute(
            sql.SQL("SELECT 1 FROM {} WHERE filename = %s LIMIT 1").format(sql.Identifier(collection_name)),
            (filename,),
        ).fetchone()

    # Return True if a result is found, False otherwise
    return row is not None


def retrieve_from_query(query, collection_name, limit=100):

    # Validation checks
    if query is None:
        raise ValueError("No query provided.")

    query_embedding = Vector(emb_text(query, openai_client, embedding_model))

    with pg_pool.connection() as conn:
        if not _collection_exists(conn, collection_name):
            raise ValueError("The collection does not exist.")

        # Search the database, ordering by cosine distance (smallest = most similar)
        rows = conn.execute(
            sql.SQL("SELECT id, filename, file_path, text, page, vector <=> %s AS distance " "FROM {} ORDER BY vector <=> %s LIMIT %s").format(
                sql.Identifier(collection_name)
            ),
            (query_embedding, query_embedding, limit),
        ).fetchall()

    # Shape the results to match the structure the callers expect (entity sub-dict)
    return [
        {
            "id": row[0],
            "distance": row[5],
            "entity": {"id": row[0], "filename": row[1], "file_path": row[2], "text": row[3], "page": row[4]},
        }
        for row in rows
    ]


#################
### DB REMOVE ###
#################


def remove_collection_from_vector_db(collection_name):
    """
    Remove an entire collection from the vector database.
    """
    with pg_pool.connection() as conn:
        if _collection_exists(conn, collection_name):
            conn.execute(sql.SQL("DROP TABLE {}").format(sql.Identifier(collection_name)))
            print(f"Collection '{collection_name}' has been removed from the vector database.")
        else:
            print(f"Collection '{collection_name}' does not exist.")
